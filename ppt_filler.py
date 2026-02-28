"""
PowerPoint Template Filler
Reads and fills the Idea Submission PowerPoint template with provided data
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import json
import sys

class PPTFiller:
    def __init__(self, template_path):
        self.template_path = template_path
        self.prs = Presentation(template_path)
    
    def analyze_template(self):
        """Analyze the template structure and print all text placeholders"""
        print(f"Analyzing template: {self.template_path}")
        print(f"Total slides: {len(self.prs.slides)}\n")
        
        for slide_idx, slide in enumerate(self.prs.slides):
            print(f"--- Slide {slide_idx + 1} ---")
            print(f"Layout: {slide.slide_layout.name}")
            
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "text"):
                    print(f"  Shape {shape_idx}: {shape.shape_type}")
                    print(f"    Text: {shape.text[:100]}")
                    
                    # Check if it has a text frame
                    if hasattr(shape, "text_frame"):
                        for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                            if paragraph.text.strip():
                                print(f"      Paragraph {para_idx}: {paragraph.text[:80]}")
            print()
    
    def fill_template(self, data):
        """Fill the template with provided data"""
        print("Filling template with data...")
        
        for slide_idx, slide_data in enumerate(data.get('slides', [])):
            if slide_idx >= len(self.prs.slides):
                print(f"Warning: Data provided for slide {slide_idx + 1} but template only has {len(self.prs.slides)} slides")
                break
            
            slide = self.prs.slides[slide_idx]
            print(f"Processing Slide {slide_idx + 1}...")
            
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    # Replace placeholders in the shape
                    for placeholder_key, replacement_value in slide_data.items():
                        if placeholder_key in shape.text:
                            shape.text = shape.text.replace(placeholder_key, str(replacement_value))
                            print(f"  Replaced '{placeholder_key}' with '{replacement_value}'")
    
    def save(self, output_path):
        """Save the modified presentation"""
        self.prs.save(output_path)
        print(f"\nPresentation saved to: {output_path}")

def main():
    template_path = "Docs/Idea Submission _ AWS AI for Bharat Hackathon.pptx"
    
    # Initialize the filler
    filler = PPTFiller(template_path)
    
    # First, analyze the template to see its structure
    print("=" * 60)
    print("TEMPLATE ANALYSIS")
    print("=" * 60)
    filler.analyze_template()
    
    print("\n" + "=" * 60)
    print("To fill the template, create a data.json file with your content")
    print("=" * 60)

if __name__ == "__main__":
    main()
